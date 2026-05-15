import pytest
from httpx import AsyncClient
from app.main import app
from reportlab.pdfgen import canvas
from pptx import Presentation
from pathlib import Path
import tempfile
import os

@pytest.fixture
async def async_client():
    async with AsyncClient(app=app, base_url="http://test") as client:
        yield client

@pytest.fixture
def sample_pdf():
    fd, path = tempfile.mkstemp(suffix=".pdf")
    os.close(fd)
    c = canvas.Canvas(path)
    c.setFont("Helvetica-Bold", 16)
    c.drawString(100, 750, "Photosynthesis")
    c.setFont("Helvetica", 12)
    c.drawString(100, 730, "Photosynthesis is a process used by plants and other organisms to convert light energy into chemical energy.")
    c.save()
    yield Path(path)
    if os.path.exists(path):
        os.remove(path)

@pytest.fixture
def sample_pptx():
    fd, path = tempfile.mkstemp(suffix=".pptx")
    os.close(fd)
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    title.text = "Water Cycle"
    subtitle = slide.placeholders[1]
    subtitle.text = "Evaporation, Condensation, Precipitation"
    
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = "Discuss the importance of solar energy in the water cycle."
    
    prs.save(path)
    yield Path(path)
    if os.path.exists(path):
        os.remove(path)
